#!/usr/bin/env python
# -*- coding: utf-8 -*-

from __future__ import (absolute_import, division,
                        print_function, unicode_literals)
from builtins import *

import re
import codecs
from txttk.retools import *
import argparse
from pptx import Presentation

def load_bible_text(path='../data/hb5.txt'):
    """
    Load the bible_text data
    """
    with codecs.open(path, encoding='big5', errors='ignore') as f:
        return f.read().strip()

def parse_bible(bible_text):
    """
    Parse the bible_text data
    """
    pharses = bible_text.split('\n')[1:]
    return pharses

def build_repository(pharses):
    """
    Build the bible repository as a dict from identifier to context
    Jhn 3:10 --> text in John chapter 3 pharse 10
    """
    repository = {}
    for pharse in pharses:
        book, _, other = pharse.partition(' ')
        locator, _, context = other.partition(' ')
        repository[' '.join([book, locator])] = context
    return repository

pharses = parse_bible(load_bible_text())
repository = build_repository(pharses)

book_map = {}
bookid2chinese = {}
with open('../data/book_names.txt') as f:
    for row in filter(lambda x: len(x)>3 and x[:4]!='中文卷名', f):
        chinese_long, chinese_short, engl_long, engl_short = row.strip().split('\t')
        the_map = {
            chinese_long: engl_short,
            chinese_short: engl_short,
            #engl_long: engl_short,
            #engl_short: engl_short
        }
        book_map.update(the_map)
        bookid2chinese[engl_short] = chinese_long


book_names = book_map.keys()
bookname_regex = condense(book_names).replace('\\', '')
locator_pattern = r'(?P<locator>[一二三四五六七八九十廿卅百\d:\-\,]+)'
filter_pattern = concat(['(?P<book>{})'.format(bookname_regex), locator_pattern])

def name_normalize(name):
    return book_map[name]

def candidate_filter(context):
    """
    Return the candidate from context
    """
    pattern = '((?P<book>加(拉太書)?|士(師記)?|(耶利米哀|雅)?歌|約(翰福音|翰貳書|翰壹書|書亞記|翰參書|伯記|珥書|拿書|三|二|一)?|耶(利米書)?|啟(示錄)?|申(命記)?|(帖撒羅尼迦前|帖撒羅尼迦後|提摩太後|撒迦利亞|提摩太前|哥林多前|俄巴底亞|哥林多後|哈巴谷|彼得後|以西結|彼得前|西番雅|阿摩司|腓利門|腓利比|歌羅西|瑪拉基|以弗所|但以理|以賽亞|希伯來|何西阿|雅各|哈該|羅馬|那鴻|彌迦|猶大|提多|傳道)?書|路(加福音|得記)?|(使徒行)?傳|民(數記)?|利(未記)?|創(世記)?|尼(希米記)?|詩(篇)?|箴(言)?|出(埃及記)?|彌|得|伯|提前|來|撒母耳記下|以斯拉記|但|何|太|拉|腓|以斯帖記|歷代志上|歷代志下|王上|徒|賽|羅|彼前|代上|門|該|拿|林後|彼後|番|帖前|撒上|撒母耳記上|哈|亞|摩|斯|帖後|俄|鴻|代下|撒下|列王紀上|馬太福音|馬可福音|列王紀下|弗|猶|王下|瑪|多|提後|可|哀|雅|結|西|珥|林前))(?P<locator>[一二三四五六七八九十廿卅百\\d:\\-\\,]+)'
    for m in re.finditer(pattern, context):
        yield m

number_map = {'一':1, 
            '二':2, 
            '三':3, 
            '四':4, 
            '五':5, 
            '六':6, 
            '七':7, 
            '八':8, 
            '九':9, 
            '十':10, 
            '廿':20, 
            '卅':30, 
            '百':100, }

def translate_number(chinese_number): 
    """
    Translate chinese number less than 999 into int
    
    >>> translate_number('一百八十二')
    182
    >>> translate_number('廿一')
    21
    """
    pattern = r'((?P<hundred>.百)?)(?P<lesser_than_100>{})'.format('[一二三四五六七八九十廿卅]*')
    number = 0
    m = re.match(pattern, chinese_number)

    if m.group('hundred'):
        number += number_map[m.group('hundred')[0]] * 100
    lesser_than_100 = m.group('lesser_than_100')
    pattern = r'(?P<ten>.十)?(?P<one>.+)'
    m = re.match(pattern, lesser_than_100)
    if m.group('ten'):
        number += number_map[m.group('ten')[0]] * 10
    if m.group('one'):
        number += sum([number_map[char] for char in m.group('one')])
    return number

def spliter(re_match_from_filter):
    """
    Split given re.match object from candidate filter into (book, locator) tuple 
    """
    book = re_match_from_filter.group('book')
    locator = re_match_from_filter.group('locator')
    return (book, locator)

def parse_range(range_text):
    """
    Split a range text such as '3-5' into [3, 4, 5]
    """
    
    sep = '-'
    if sep in range_text:
        start, end = (int(n) for n in range_text.split(sep))
        return list(range(start, end + 1))
    else:
        return [int(range_text)]
    
def parse_locator(locator):
    """
    Parse given locator such as '一5', '5章2-3'
    """
    pattern = r'(?:(?P<number>\d+)|(?P<chinese>[一二三四五六七八九十廿卅百]+))[章:]?(?P<pharse_range>[\d,\-]+)'
    m = re.match(pattern, locator)
    if m.group('number'):
        chapter = int(m.group('number'))
    else:
        chapter = translate_number(m.group('chinese'))
    range_text = m.group('pharse_range')
    pharse_list = parse_range(range_text)
    return chapter, pharse_list

def get_bucket(re_match_from_filter):
    """
    From given re.match object from candidate_filter, 
    return a bucket (a list) of bible identifiers (bookid, chapter, pharse)
    """
    book, locator = spliter(re_match_from_filter)
    chapter, pharse_list = parse_locator(locator)
    bucket = [(name_normalize(book), chapter, pharse) for pharse in pharse_list]
    return bucket

def get_context(book, chapter, pharse):
    """
    Given book, chapter, and pharse number, return the bible context.
    """
    context = repository['{} {}:{}'.format(book, chapter, pharse)]
    return context

def format_bucket(bucket):
    bookids = [pharse[0] for pharse in bucket]
    assert len(set(bookids)) == 1
    bookid = bookids[0]
    bookname = bookid2chinese[bookid]
    chapters = [pharse[1] for pharse in bucket]
    assert len(set(chapters)) == 1
    chapter = chapters[0]
    header = '{}{}章'.format(bookname, chapter)
    body = ''.join(['{}{}'.format(p[2], get_context(*p)) for p in bucket])
    return header, body
    
def to_pptx(filename, pages):
    """
    A page is a tuple of (title, text)
    """
    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]
    for page in pages:
        
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = page[0]

        tf = body_shape.text_frame
        tf.text = page[1]

    prs.save(filename)
    
def open_input(filename):
    """
    Read the inputfile, try big5 and utf8 codecs
    """
    try:
        with codecs.open(filename, mode='r', encoding='utf8') as f:
            text = f.read()
    except UnicodeDecodeError:
        with codecs.open(filename, mode='r', encoding='big5') as f:
            text = f.read()
    return text
        
def main(args):
    input_filename = args.input
    output_filename = args.output
    
    context = open_input(input_filename)
    pages = []
    for candidate in candidate_filter(context):
        bucket = get_bucket(candidate)
        page = format_bucket(bucket)
        pages.append(page)
    to_pptx(output_filename, pages)



########################


if __name__ == '__main__':
    parser = argparse.ArgumentParser( 
        description = "From bible locators in a text file to generate a powerpoint file contains all the scripture.",
        epilog = "As an alternative to the commandline, params can be placed in a file, one per line, and specified on the commandline like '%(prog)s @params.conf'.",
        fromfile_prefix_chars = '@' )
    # TODO Specify your real parameters here.
    parser.add_argument(
                      "-i",
                      "--input",
                      help="the input text file",
                      action="store")
    parser.add_argument(
                      "-o",
                      "--output",
                      help = "the output file ends with pptx",
                      action="store")
    
    args = parser.parse_args()

    main(args)
